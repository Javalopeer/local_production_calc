"""Generate a professional PowerPoint presentation for ProductionCalcApp."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "ProductionCalcApp_Presentation.pptx",
)

# ── Brand colors ──────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x38, 0x64)
MED_BLUE    = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
ACCENT_GREEN= RGBColor(0x54, 0x82, 0x35)
ACCENT_RED  = RGBColor(0xC6, 0x28, 0x28)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
BLACK       = RGBColor(0x00, 0x00, 0x00)


def _add_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_slide(prs, title_text, subtitle_text, bullets, color_accent=MED_BLUE):
    """Create a slide with title, subtitle, and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, WHITE)

    # Top accent bar
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), color_accent)

    # Title
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                  title_text, font_size=32, color=DARK_BLUE, bold=True)

    # Subtitle
    if subtitle_text:
        _add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                      subtitle_text, font_size=16, color=MED_BLUE, bold=False)

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(bullet, tuple):
            # (main_text, sub_text)
            p.text = f"▸  {bullet[0]}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.font.bold = True
            p.font.name = "Calibri"
            p.space_after = Pt(2)
            # Sub-bullet
            p2 = tf.add_paragraph()
            p2.text = f"      {bullet[1]}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.font.name = "Calibri"
            p2.space_after = Pt(12)
        else:
            p.text = f"▸  {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.font.bold = False
            p.font.name = "Calibri"
            p.space_after = Pt(10)

    return slide


def _add_two_column_slide(prs, title_text, left_items, right_items,
                           left_title="", right_title=""):
    """Create a two-column slide with bullet lists."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), MED_BLUE)

    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                  title_text, font_size=32, color=DARK_BLUE, bold=True)

    # Left column title
    if left_title:
        _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
                      left_title, font_size=18, color=MED_BLUE, bold=True)

    # Left bullets
    txL = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0))
    tf = txL.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(8)

    # Right column title
    if right_title:
        _add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.4),
                      right_title, font_size=18, color=MED_BLUE, bold=True)

    # Right bullets
    txR = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.0))
    tf2 = txR.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(8)

    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BLUE)

    # Accent bar
    _add_shape(slide, Inches(0), Inches(3.0), Inches(13.33), Inches(0.06), MED_BLUE)

    _add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
                  "ProductionCalcApp", font_size=48, color=WHITE, bold=True)

    _add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.6),
                  "Production Performance Calculator", font_size=24, color=LIGHT_BLUE)

    _add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                  "Aplicación de escritorio para tracking, análisis y optimización\n"
                  "de la producción del equipo de diseño",
                  font_size=18, color=RGBColor(0xAA, 0xBB, 0xDD))

    _add_text_box(slide, Inches(1), Inches(5.5), Inches(5), Inches(0.4),
                  "Desarrollado por: Gerardo", font_size=16, color=LIGHT_BLUE)
    _add_text_box(slide, Inches(1), Inches(5.9), Inches(5), Inches(0.4),
                  "Versión 1.1.4  •  Marzo 2026", font_size=14,
                  color=RGBColor(0x88, 0x99, 0xBB))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: Agenda / Contenido
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Agenda", "", [
        "Visión general de la aplicación",
        "Funcionalidades principales",
        "Sistema de registro de casos",
        "Dashboard y métricas en tiempo real",
        "Gestión de tiempos muertos (Downtimes)",
        "Sistema de aprobación de downtimes",
        "Performance tracking y justificaciones",
        "Exportación y sync con SharePoint/OneDrive",
        "Cambios recientes y mejoras",
        "Próximos pasos",
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Vision General
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Visión General", "¿Qué es ProductionCalcApp?", [
        ("Aplicación de escritorio para Windows",
         "Desarrollada en Python con PySide6/Qt — interfaz moderna y responsiva"),
        ("Control de producción individual",
         "Cada diseñador registra sus casos, tiempos y métricas en su propia instancia"),
        ("Base de datos local con sync automático",
         "SQLite local sincronizada vía OneDrive — datos seguros y accesibles"),
        ("Reportes compartidos en tiempo real",
         "Exportación automática a Excel en SharePoint para visibilidad del equipo"),
        ("15+ diseñadores en producción",
         "Desplegada y en uso activo por todo el equipo de diseño"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Arquitectura Técnica
    # ══════════════════════════════════════════════════════════════════════════
    _add_two_column_slide(prs, "Arquitectura Técnica",
        left_items=[
            "Python 3.13 + PySide6 (Qt)",
            "SQLite — base de datos local",
            "openpyxl — generación de Excel",
            "PyInstaller — empaquetado .exe",
            "qtawesome — iconos Font Awesome",
            "OneDrive — sincronización de archivos",
        ],
        right_items=[
            "~10,000+ líneas de código",
            "25+ archivos Python",
            "5 tablas en base de datos",
            "Migraciones automáticas de schema",
            "Auto-installer incluido",
            ".exe de 60 MB, sin dependencias",
        ],
        left_title="Tecnologías",
        right_title="Métricas del Proyecto",
    )

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Register Tab
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Registro de Casos", "Register Tab — El corazón de la app", [
        ("Registro rápido de casos",
         "Case ID, región, tipo de caso, doctor, hora inicio/fin — todo en un formulario intuitivo"),
        ("Importación automática desde clipboard",
         "Copiar datos de otras fuentes y pegar directamente — auto-fill inteligente"),
        ("Importación desde web (Case Scraper)",
         "Importar datos del caso directamente desde el sistema web con un clic"),
        ("Cálculo automático de tiempo real",
         "Diferencia entre hora inicio y fin calculada automáticamente"),
        ("Cálculo de case value y eficiencia",
         "Comparación con tiempos estándar para medir rendimiento por caso"),
        ("Edición y eliminación de casos",
         "Modificar o borrar casos ya registrados sin complicaciones"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Production Tab
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Vista de Producción", "Production Tab — Métricas del día", [
        ("Tabla de casos del día con filtros",
         "Filtrar por fecha, ver todos los casos registrados con detalles completos"),
        ("Cálculo de unidades equivalentes (UE)",
         "Motor de cálculo basado en estándares configurables por región y tipo"),
        ("Porcentaje de producción en tiempo real",
         "Actualización inmediata al registrar cada caso"),
        ("Paginación para manejo de muchos casos",
         "Navegación eficiente cuando hay muchos registros"),
        ("Columnas dinámicas adaptables",
         "La tabla se ajusta según el contenido y tamaño de ventana"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Dashboard
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Dashboard Interactivo",
                      "Dashboard Tab — Vista ejecutiva de producción", [
        ("KPI Cards en tiempo real",
         "Producción %, UE, casos totales, downtime — métricas clave de un vistazo"),
        ("Gráficas de producción",
         "Visualización de tendencias y distribución por región/tipo de caso"),
        ("Resumen diario consolidado",
         "Integra casos regulares, overtime y downtimes en una vista unificada"),
        ("Actualización automática",
         "Se refresca al guardar cada caso — siempre actualizado"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Overtime Tab
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Gestión de Overtime", "Overtime Tab — Casos fuera de horario", [
        ("Registro separado de casos OT",
         "Los casos después de las 3:00 PM se registran aparte para tracking diferenciado"),
        ("Misma funcionalidad que Register",
         "Importación clipboard, case scraper, cálculo de UE — todo disponible"),
        ("Producción OT independiente",
         "Métricas de overtime no afectan el target de producción regular"),
        ("Historial de OT por día",
         "Visualización de carga de trabajo extra por fecha"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9: History Tab
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Historial de Casos", "History Tab — Registro completo", [
        ("Historial completo de todos los casos",
         "Búsqueda y filtrado por fecha, tipo, región — acceso rápido a cualquier registro"),
        ("Filtros avanzados",
         "Combinación de filtros por rango de fechas, tipo de caso y región"),
        ("Exportación a Excel",
         "Exportar historial filtrado a archivo Excel para análisis externo"),
        ("Estadísticas de resumen",
         "Totales, promedios y métricas agregadas del periodo filtrado"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10: Downtime Management
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Gestión de Tiempos Muertos",
                      "Downtime Manager — Control de tiempo no productivo", [
        ("Registro de downtimes con categorización",
         "Tipo de downtime, duración, motivo — todo clasificado para análisis"),
        ("Impacto automático en producción",
         "Los downtimes aprobados se restan del tiempo base, ajustando el % de producción"),
        ("Cálculo de UE por downtime",
         "Los downtimes generan unidades equivalentes proporcionales al tiempo"),
        ("Sistema de aprobación por supervisor",
         "Los downtimes requieren aprobación — Excel compartido en OneDrive"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 11: Approval System
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Sistema de Aprobación de Downtimes",
                      "Workflow automático multi-usuario", [
        ("Excel compartido en OneDrive",
         "Archivo _Downtime_Approvals.xlsx en carpeta Downtime/ — accesible por supervisores"),
        ("Protección de hojas",
         "Solo la columna Status es editable — el resto está protegido con contraseña"),
        ("Polling automático cada 15 segundos",
         "La app consulta periódicamente si el supervisor aprobó/rechazó downtimes"),
        ("Botón manual 'Check Approvals'",
         "El diseñador puede verificar manualmente el estado de sus solicitudes"),
        ("Actualización automática de DB",
         "Al detectar aprobación/rechazo, la base de datos local se actualiza inmediatamente"),
        ("Multi-diseñador",
         "Cada diseñador solo ve y modifica sus propios downtimes en el archivo compartido"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 12: Performance Tracking
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Performance Tracking",
                      "Sistema de seguimiento diario de rendimiento", [
        ("Popup de felicitación al alcanzar el target",
         "Cuando el diseñador alcanza ≥95% de producción o ≥14 UE, se muestra una felicitación"),
        ("Popup de justificación al no alcanzar el target",
         "A las 3:15 PM, si no se alcanzó el target, se requiere una justificación escrita"),
        ("No se puede cerrar la app sin justificar",
         "Después de las 3:15 PM, la app bloquea el cierre hasta enviar la justificación"),
        ("Enforcement al día siguiente",
         "Si el diseñador no justificó, al abrir la app al siguiente día se le pide antes de poder trabajar"),
        ("Casos sobre estándar incluidos",
         "La justificación incluye automáticamente los casos que excedieron el tiempo estándar"),
        ("Exportación a Excel compartido",
         "Las justificaciones se guardan en _Daily_Justifications.xlsx con filtros por diseñador y fecha"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 13: SharePoint Sync
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Exportación y Sincronización",
                      "SharePoint / OneDrive Integration", [
        ("Sync automático al guardar cada caso",
         "Al registrar un caso, se exporta automáticamente al Excel de SharePoint"),
        ("Dashboard compartido en Excel",
         "Archivo consolidado con producción de todos los diseñadores — visible para supervisores"),
        ("Columnas separadas de UE",
         "UE (Cases) y UE (w/ DT) — producción pura vs producción con downtime"),
        ("Auto-detección de carpeta OneDrive",
         "La app detecta automáticamente la ruta de OneDrive sin configuración manual"),
        ("Archivos protegidos",
         "Los Excel compartidos tienen protección con contraseña para evitar edición accidental"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 14: Standards Tab
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Estándares Configurables",
                      "Standards Tab — Tiempos estándar por región y tipo", [
        ("Gestión de tiempos estándar",
         "Cada región y tipo de caso tiene un tiempo estándar configurable"),
        ("Import/Export JSON",
         "Los estándares se pueden importar y exportar para compartir configuraciones"),
        ("Impacto inmediato en cálculos",
         "Al modificar un estándar, todos los cálculos de UE se actualizan automáticamente"),
        ("Interfaz de árbol jerárquico",
         "Vista organizada por región → tipo de caso → tiempo estándar"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 15: UI/UX
    # ══════════════════════════════════════════════════════════════════════════
    _add_two_column_slide(prs, "Experiencia de Usuario",
        left_items=[
            "Tema oscuro y claro con toggle",
            "Interfaz responsiva y adaptable",
            "Iconos Font Awesome profesionales",
            "Tablas con colores alternados",
            "Tooltips informativos",
            "Atajos de teclado",
        ],
        right_items=[
            "Toast notifications al importar",
            "Diálogos de confirmación para acciones destructivas",
            "Progress bars para operaciones largas",
            "Auto-fill inteligente de formularios",
            "Scroll areas para contenido extenso",
            "Freeze panes en Excel exportados",
        ],
        left_title="Diseño Visual",
        right_title="Interacción",
    )

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 16: Cambios Recientes
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Cambios Recientes",
                      "Actualizaciones Marzo 2026", [
        ("Separación de columnas UE en dashboard",
         "UE (Cases) y UE (w/ DT) para diferenciar producción pura vs con downtime"),
        ("Sistema completo de aprobación de downtimes",
         "Excel compartido, protección de hojas, polling automático, multi-diseñador"),
        ("Performance tracking con popups",
         "Felicitación al 95%, justificación requerida al no alcanzar, enforcement al siguiente día"),
        ("Casos sobre estándar en justificaciones",
         "Al justificar, se incluyen automáticamente los casos que excedieron el tiempo permitido"),
        ("Auto-filtro en Excel compartidos",
         "Filtrar por diseñador y fecha directamente en Excel — protección permite filtrar y ordenar"),
        ("OneDrive force-refresh",
         "Mecanismo para forzar la descarga de la última versión del archivo de aprobaciones"),
        ("Legacy DB migration",
         "Migración automática de bases de datos existentes al nuevo schema"),
        ("Build .exe optimizado",
         "Ejecutable de 60 MB con auto-installer — distribución sin dependencias"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 17: Deployment
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Deployment y Distribución",
                      "Instalación y actualización simplificada", [
        ("Ejecutable standalone (.exe)",
         "Un solo archivo de 60 MB — no requiere Python ni dependencias instaladas"),
        ("Auto-installer incluido",
         "Al ejecutar por primera vez, se instala automáticamente en la ubicación correcta"),
        ("Base de datos en OneDrive",
         "Los datos se almacenan en %OneDrive%/ProductionCalcApp/ — backup automático en la nube"),
        ("Actualización sin pérdida de datos",
         "Al actualizar el .exe, todos los casos, downtimes y configuraciones se mantienen intactos"),
        ("Schema migrations automáticas",
         "CREATE TABLE IF NOT EXISTS y ALTER TABLE ADD COLUMN con error handling"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 18: Métricas del Proyecto
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), MED_BLUE)

    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                  "Métricas del Proyecto", font_size=32, color=DARK_BLUE, bold=True)

    # KPI Cards
    kpis = [
        ("85+", "Horas de\nDesarrollo", MED_BLUE),
        ("87", "Commits\nRegistrados", ACCENT_GREEN),
        ("13,800+", "Líneas de\nCódigo", ORANGE),
        ("15+", "Diseñadores\nen Producción", DARK_BLUE),
        ("25+", "Archivos\nPython", MED_BLUE),
        ("19", "Días de\nTrabajo", ACCENT_GREEN),
    ]
    card_width = Inches(1.7)
    card_height = Inches(1.8)
    start_x = Inches(0.8)
    gap = Inches(0.25)

    for i, (value, label, color) in enumerate(kpis):
        x = start_x + i * (card_width + gap)
        y = Inches(1.5)

        # Card background
        card = _add_shape(slide, x, y, card_width, card_height, LIGHT_GRAY)

        # Top color bar
        _add_shape(slide, x, y, card_width, Inches(0.06), color)

        # Value
        _add_text_box(slide, x, y + Inches(0.2), card_width, Inches(0.7),
                      value, font_size=36, color=color, bold=True,
                      alignment=PP_ALIGN.CENTER)

        # Label
        _add_text_box(slide, x, y + Inches(0.95), card_width, Inches(0.7),
                      label, font_size=13, color=DARK_GRAY,
                      alignment=PP_ALIGN.CENTER)

    # Timeline
    _add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.4),
                  "Línea de Tiempo del Desarrollo", font_size=20,
                  color=DARK_BLUE, bold=True)

    timeline = [
        ("Ene 27-29", "Estructura base,\nDB, tabs principales"),
        ("Feb 4-6", "UE calculation,\ntemas, responsive"),
        ("Feb 13", "Unit tests,\nstyling cross-theme"),
        ("Feb 23", "Dashboard, sync,\nPyInstaller"),
        ("Mar 3", "SharePoint export,\nauto-config"),
        ("Mar 17-20", "Approvals, performance\ntracking, popups"),
    ]
    for i, (period, desc) in enumerate(timeline):
        x = Inches(0.8) + i * Inches(2.0)
        y = Inches(4.4)

        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y,
                                      Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = MED_BLUE
        dot.line.fill.background()

        # Period
        _add_text_box(slide, x, y + Inches(0.3), Inches(1.8), Inches(0.3),
                      period, font_size=12, color=MED_BLUE, bold=True,
                      alignment=PP_ALIGN.CENTER)

        # Description
        _add_text_box(slide, x, y + Inches(0.6), Inches(1.8), Inches(0.8),
                      desc, font_size=10, color=DARK_GRAY,
                      alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 19: Próximos Pasos
    # ══════════════════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Próximos Pasos", "Roadmap de funcionalidades futuras", [
        ("Reportes semanales/mensuales automáticos",
         "PDF o Excel con resumen de producción por diseñador, enviado al supervisor"),
        ("Modo supervisor",
         "Vista especial para supervisores con métricas de todo el equipo y aprobación directa"),
        ("Gráficas de tendencia personal",
         "Cada diseñador ve su progreso semanal/mensual con línea de tendencia"),
        ("Ranking de equipo",
         "Vista comparativa del desempeño del equipo para motivar competencia sana"),
        ("Predicción de fin de día",
         "Basado en el ritmo actual, estimar si el diseñador cumplirá el target"),
        ("Categorización de downtimes",
         "Reportes de cuáles son los downtimes más frecuentes para mejorar procesos"),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 20: Thank You
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BLUE)
    _add_shape(slide, Inches(0), Inches(3.2), Inches(13.33), Inches(0.06), MED_BLUE)

    _add_text_box(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(1.0),
                  "¡Gracias!", font_size=52, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(3.6), Inches(11.33), Inches(0.6),
                  "ProductionCalcApp v1.1.4", font_size=22, color=LIGHT_BLUE,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(4.3), Inches(11.33), Inches(0.6),
                  "¿Preguntas o sugerencias?", font_size=18,
                  color=RGBColor(0xAA, 0xBB, 0xDD),
                  alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f"\n✅ Presentación generada: {OUTPUT}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
